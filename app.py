import streamlit as st
from pptx import Presentation
from pptx.util import Pt
import re
import io
from PIL import Image
import base64

# 페이지 설정
st.set_page_config(layout="wide")

# 사용자 정의 CSS 추가 (최상위 커스터마이징)
st.markdown(
    """
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Pretendard&display=swap');
        html, body, [class*="css"]  {
            font-family: 'Pretendard', sans-serif;
            background-color: #FAFAFA;
        }
        .title {
            font-size: 36px;
            font-weight: 700;
            margin-bottom: 20px;
            color: #333;
        }
        .subtitle {
            font-size: 18px;
            color: #666;
            margin-bottom: 40px;
        }
        .button-container button {
            transition: transform 0.2s ease;
        }
        .button-container button:active {
            transform: scale(0.96);
        }
        .loader {
            border: 5px solid #f3f3f3;
            border-top: 5px solid #5A5A5A;
            border-radius: 50%;
            width: 40px;
            height: 40px;
            animation: spin 1s linear infinite;
            margin: auto;
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
    </style>
    """,
    unsafe_allow_html=True
)

def split_text(text, min_len=80, max_len=110):
    sentences = re.split(r'(\.|\?|!|\n)', text)
    combined = []
    temp = ''

    for i in range(0, len(sentences) - 1, 2):
        sentence = (sentences[i] + sentences[i+1]).strip()
        if not sentence:
            continue
        if len(temp) + len(sentence) <= max_len:
            temp += ' ' + sentence if temp else sentence
        else:
            if len(temp) >= min_len:
                combined.append(temp.strip())
                temp = sentence
            else:
                temp += ' ' + sentence

    if temp:
        combined.append(temp.strip())
    
    return combined

import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
import io

# 제목
st.title("Prompt to PPT")

# 입력창
text_input = st.text_area("프롬프트를 입력하세요", height=300)

# 파일 업로드
uploaded_file = st.file_uploader("PPT 양식 파일을 업로드하세요", type=["pptx"])

# 슬라이드 구분
def split_text(text):
    return [slide.strip() for slide in text.split("\n") if slide.strip()]

# PPT 생성
def create_ppt(slides, filename):
    prs = Presentation(filename)
    blank_slide_layout = prs.slide_layouts[6]  # 완전 빈 슬라이드 (Title Only가 아님)

    for slide_text in slides:
        slide = prs.slides.add_slide(blank_slide_layout)

        # 텍스트박스 새로 추가
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = slide_text
        p.font.size = Pt(40)  # 40pt
        p.font.name = '맑은 고딕'  # 맑은 고딕
        p.alignment = 1  # 가운데 정렬

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# 변환 버튼
if st.button("PPT 변환하기"):
    if text_input and uploaded_file:
        slides = split_text(text_input)
        ppt_file = create_ppt(slides, uploaded_file)
        st.download_button(
            label="PPT 다운로드",
            data=ppt_file,
            file_name="converted_ppt.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    else:
        st.error("프롬프트와 PPT 양식 파일을 모두 입력하세요.")

def estimate_time(text):
    text_no_space = re.sub(r'\s+', '', text)
    total_chars = len(text_no_space)
    minutes = total_chars / 250
    seconds = minutes * 60
    return int(seconds // 60), int(seconds % 60)

# Streamlit 앱 시작
st.markdown("<div class='title'>📂강의용 프롬프트 자동 변환하기</div>", unsafe_allow_html=True)

st.write("""
☁️ 텍스트를 입력하면, 강의용 프롬프트로 변환해드립니다.  
💬 **한 슬라이드에 100자 이내** 작성을 추천드립니다.  
""")

filename = st.text_input("📥 파일명을 입력해 주세요.", "[프롬프트] 과정명_O차시_이름_날짜")
char_limit = st.number_input("✨ 한 슬라이드에 들어갈 글자수를 입력해 주세요.", min_value=50, max_value=150, value=100, step=10)
prompt = st.text_area("💭 내용을 작성해 주세요. (작성 후, Ctrl+Enter)", "")

if prompt:
    slides = split_text(prompt, min_len=80, max_len=char_limit)

    if st.button("🔍 문장 슬라이드 미리보기"):
        st.subheader("📝 슬라이드 미리보기 ")
        for i, slide in enumerate(slides):
            st.markdown(f"☁️ 슬라이드 {i+1} ")
            st.code(slide, language='text')

    if st.button("⏰ 예상 발표시간 확인하기"):
        minutes, seconds = estimate_time(prompt)
        st.success(f"⏳ 예상 발표시간: 약 {minutes}분 {seconds}초")
        if minutes < 20:
            st.info("📢 차시당 20~25분 분량을 권장드립니다.")

    if st.button("✅ 이대로 PPT로 변환하기"):
        ppt_file = create_ppt(slides, filename)
        st.download_button(
            label="📥 다운로드",
          data=ppt_file,
          file_name=f"{filename}.pptx",
         mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    st.toast('PPT 파일이 성공적으로 생성되었습니다!', icon='🎉')

    # 페이지 하단에 고정된 문의사항 문구 추가
footer = """
    <style>
        .footer {
            position: fixed;
            bottom: 10px;
            right: 10px;
            font-size: 0.9em;
            color: #666;
            z-index: 100;
        }
    </style>
    <div class="footer">
        문의사항: kimhj11@visang.com
    </div>
"""
st.markdown(footer, unsafe_allow_html=True)

